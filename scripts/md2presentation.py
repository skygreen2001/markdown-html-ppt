#!/usr/bin/env python3
"""
Markdown Html Presentation Converter v2.4.0
将 Markdown 文档转换为演示文稿（HTML reveal.js 或 PowerPoint）
支持内容自动分页、多主题和丰富的 Markdown 语法
修复：HTML 标签正确闭合
"""

import re
import argparse
import sys
from pathlib import Path
from typing import List, Tuple, Optional, Union
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor

# 导入主题模块
try:
    from themes import get_theme, Theme, list_themes
except ImportError:
    print("警告: 无法导入主题模块，将使用默认主题")
    Theme = None
    get_theme = None
    list_themes = None


class EnhancedMarkdownParser:
    """增强的 Markdown 解析器，支持更多语法"""
    
    @staticmethod
    def parse_inline_formatting(text: str) -> List[dict]:
        """
        解析行内格式（加粗、斜体、代码）
        
        Returns:
            格式化文本片段列表
        """
        segments = []
        current_pos = 0
        
        # 正则表达式匹配加粗、斜体、代码
        pattern = r'(\*\*|__)(.*?)\1|(\*|_)(.*?)\3|(`)(.*?)\5'
        
        for match in re.finditer(pattern, text):
            # 添加匹配前的普通文本
            if match.start() > current_pos:
                plain_text = text[current_pos:match.start()]
                if plain_text:
                    segments.append({
                        'text': plain_text,
                        'bold': False,
                        'italic': False,
                        'code': False
                    })
            
            # 判断格式类型
            if match.group(1) in ('**', '__'):  # 加粗
                segments.append({
                    'text': match.group(2),
                    'bold': True,
                    'italic': False,
                    'code': False
                })
            elif match.group(3) in ('*', '_'):  # 斜体
                segments.append({
                    'text': match.group(4),
                    'bold': False,
                    'italic': True,
                    'code': False
                })
            elif match.group(5) == '`':  # 代码
                segments.append({
                    'text': match.group(6),
                    'bold': False,
                    'italic': False,
                    'code': True
                })
            
            current_pos = match.end()
        
        # 添加剩余的普通文本
        if current_pos < len(text):
            remaining = text[current_pos:]
            if remaining:
                segments.append({
                    'text': remaining,
                    'bold': False,
                    'italic': False,
                    'code': False
                })
        
        # 如果没有任何格式，返回整个文本
        if not segments:
            segments.append({
                'text': text,
                'bold': False,
                'italic': False,
                'code': False
            })
        
        return segments
    
    @staticmethod
    def parse(content: str) -> List[dict]:
        """
        解析 Markdown 内容为结构化数据
        
        Args:
            content: Markdown 文本内容
            
        Returns:
            幻灯片数据列表
        """
        slides = []
        current_slide = None
        lines = content.split('\n')
        
        i = 0
        while i < len(lines):
            line = lines[i].rstrip()
            
            # 一级标题 - 标题幻灯片
            if line.startswith('# '):
                if current_slide:
                    slides.append(current_slide)
                current_slide = {
                    'type': 'title',
                    'title': line[2:].strip(),
                    'subtitle': '',
                    'content': []
                }
                i += 1
                continue
            
            # 二级/三级标题 - 内容幻灯片
            elif line.startswith('## ') or line.startswith('### '):
                if current_slide:
                    slides.append(current_slide)
                level = 2 if line.startswith('## ') else 3
                title = line[level+1:].strip()
                current_slide = {
                    'type': 'content',
                    'title': title,
                    'content': []
                }
                i += 1
                continue
            
            # 块引用
            elif line.startswith('> '):
                if current_slide:
                    quote_lines = []
                    while i < len(lines) and lines[i].startswith('> '):
                        quote_lines.append(lines[i][2:].strip())
                        i += 1
                    
                    quote_text = ' '.join(quote_lines)
                    current_slide['content'].append({
                        'type': 'quote',
                        'segments': EnhancedMarkdownParser.parse_inline_formatting(quote_text)
                    })
                    continue
            
            # 代码块
            elif line.strip().startswith('```'):
                if current_slide:
                    i += 1
                    code_lines = []
                    while i < len(lines) and not lines[i].strip().startswith('```'):
                        code_lines.append(lines[i])
                        i += 1
                    
                    if i < len(lines):
                        i += 1  # 跳过结束的 ```
                    
                    current_slide['content'].append({
                        'type': 'code',
                        'text': '\n'.join(code_lines)
                    })
                    continue
            
            # 列表项
            elif line.strip().startswith(('- ', '* ', '+ ')) or re.match(r'^\s*\d+\.\s', line):
                if current_slide:
                    content_text = re.sub(r'^[\s\-\*\+\d\.]+\s*', '', line).strip()
                    if content_text:
                        current_slide['content'].append({
                            'type': 'list',
                            'segments': EnhancedMarkdownParser.parse_inline_formatting(content_text)
                        })
                i += 1
                continue
            
            # 普通段落
            elif line.strip() and current_slide:
                if current_slide['type'] == 'title' and not current_slide['subtitle']:
                    current_slide['subtitle'] = line.strip()
                elif line.strip():
                    current_slide['content'].append({
                        'type': 'paragraph',
                        'segments': EnhancedMarkdownParser.parse_inline_formatting(line.strip())
                    })
                i += 1
                continue
            
            i += 1
        
        if current_slide:
            slides.append(current_slide)
        
        return slides


class SmartPaginator:
    """智能分页器 - 根据内容实际高度计算"""
    
    # 字体大小和行高配置
    TITLE_FONT_SIZE = 28
    CONTENT_FONT_SIZE = 14
    LINE_SPACING = 1.5
    
    # 幻灯片尺寸（点）
    SLIDE_HEIGHT_PT = 405
    SLIDE_WIDTH_PT = 720
    
    # 可用内容区域
    CONTENT_TOP_PT = 80
    CONTENT_BOTTOM_PT = 380
    CONTENT_HEIGHT_PT = CONTENT_BOTTOM_PT - CONTENT_TOP_PT  # 300pt
    
    @staticmethod
    def estimate_content_height(content_items: List[dict]) -> float:
        """估算内容的实际高度（点）"""
        total_height = 0
        
        for item in content_items:
            item_type = item.get('type', 'list')
            
            if item_type == 'quote':
                # 引用：基础高度 + 15%
                segments = item.get('segments', [])
                text_length = sum(len(seg.get('text', '')) for seg in segments)
                height = (text_length / 50) * SmartPaginator.CONTENT_FONT_SIZE * SmartPaginator.LINE_SPACING
                total_height += height * 1.15 + 20  # 额外边距
            
            elif item_type == 'code':
                # 代码块：每行高度
                code_text = item.get('text', '')
                lines = code_text.count('\n') + 1
                total_height += lines * SmartPaginator.CONTENT_FONT_SIZE * 1.2 + 20
            
            elif item_type == 'paragraph':
                # 段落
                segments = item.get('segments', [])
                text_length = sum(len(seg.get('text', '')) for seg in segments)
                height = (text_length / 50) * SmartPaginator.CONTENT_FONT_SIZE * SmartPaginator.LINE_SPACING
                total_height += height + 10
            
            else:  # list
                # 列表项
                segments = item.get('segments', [])
                text_length = sum(len(seg.get('text', '')) for seg in segments)
                
                # 根据文本长度和格式调整高度
                base_height = SmartPaginator.CONTENT_FONT_SIZE * SmartPaginator.LINE_SPACING
                
                # 检查是否有格式化
                has_formatting = any(seg.get('bold') or seg.get('italic') or seg.get('code') 
                                    for seg in segments)
                if has_formatting:
                    base_height *= 1.1
                
                # 根据文本长度调整
                if text_length > 100:
                    base_height *= 1.5
                elif text_length > 50:
                    base_height *= 1.2
                
                total_height += base_height + 6
        
        return total_height
    
    @staticmethod
    def auto_paginate(slides_data: List[dict]) -> List[dict]:
        """智能自动分页"""
        paginated_slides = []
        
        for slide in slides_data:
            if slide['type'] != 'content' or not slide['content']:
                paginated_slides.append(slide)
                continue
            
            # 估算总高度
            total_height = SmartPaginator.estimate_content_height(slide['content'])
            
            # 如果高度在可用范围内，不分页
            if total_height <= SmartPaginator.CONTENT_HEIGHT_PT:
                paginated_slides.append(slide)
                continue
            
            # 需要分页：动态计算每页容量
            content_items = slide['content']
            current_page_items = []
            current_height = 0
            page_num = 0
            
            for item in content_items:
                item_height = SmartPaginator.estimate_content_height([item])
                
                # 检查是否需要新页
                if current_height + item_height > SmartPaginator.CONTENT_HEIGHT_PT and current_page_items:
                    # 创建新幻灯片
                    new_slide = {
                        'type': 'content',
                        'title': slide['title'] + (f" (续 {page_num})" if page_num > 0 else ""),
                        'content': current_page_items
                    }
                    paginated_slides.append(new_slide)
                    
                    # 重置
                    current_page_items = [item]
                    current_height = item_height
                    page_num += 1
                else:
                    current_page_items.append(item)
                    current_height += item_height
            
            # 添加最后一页
            if current_page_items:
                new_slide = {
                    'type': 'content',
                    'title': slide['title'] + (f" (续 {page_num})" if page_num > 0 else ""),
                    'content': current_page_items
                }
                paginated_slides.append(new_slide)
        
        return paginated_slides


class RevealJSGenerator:
    """reveal.js HTML 生成器 - 修复版"""
    
    @staticmethod
    def _escape_html(text: str) -> str:
        """转义 HTML 特殊字符"""
        return (text
                .replace('&', '&amp;')
                .replace('<', '&lt;')
                .replace('>', '&gt;')
                .replace('"', '&quot;')
                .replace("'", '&#39;'))
    
    @staticmethod
    def _format_inline(segments: List[dict]) -> str:
        """格式化行内文本"""
        result = ""
        for seg in segments:
            text = RevealJSGenerator._escape_html(seg.get('text', ''))
            
            if seg.get('code'):
                result += f'<code style="background: #f0f0f0; padding: 2px 6px; border-radius: 3px; font-family: monospace;">{text}</code>'
            elif seg.get('bold') and seg.get('italic'):
                result += f'<strong><em>{text}</em></strong>'
            elif seg.get('bold'):
                result += f'<strong>{text}</strong>'
            elif seg.get('italic'):
                result += f'<em>{text}</em>'
            else:
                result += text
        
        return result
    
    @staticmethod
    def generate(slides_data: List[dict], title: str = "演示文稿", theme: Theme = None) -> str:
        """
        生成 reveal.js HTML
        
        Args:
            slides_data: 幻灯片数据
            title: 演示文稿标题
            theme: 主题对象
            
        Returns:
            HTML 字符串
        """
        if theme is None and get_theme:
            theme = get_theme('default')
        
        sections = []
        
        for slide in slides_data:
            if slide['type'] == 'title':
                # 标题幻灯片
                section = f'''
            <section class="center">
                <h1>{RevealJSGenerator._escape_html(slide['title'])}</h1>'''
                if slide.get('subtitle'):
                    section += f'''
                <h3>{RevealJSGenerator._escape_html(slide['subtitle'])}</h3>'''
                section += '''
            </section>'''
            else:
                # 内容幻灯片
                section = f'''
            <section>
                <h2>{RevealJSGenerator._escape_html(slide['title'])}</h2>'''
                
                if slide['content']:
                    # 🔧 修复：添加状态跟踪，确保 ul 标签正确闭合
                    in_list = False
                    
                    for item in slide['content']:
                        item_type = item.get('type', 'list')
                        
                        # 检查是否需要关闭列表
                        if in_list and item_type != 'list':
                            section += '''
                </ul>'''
                            in_list = False
                        
                        if item_type == 'quote':
                            # 块引用
                            section += '''
                <blockquote style="font-style: italic; border-left: 3px solid #ccc; padding-left: 15px; margin: 10px 0; line-height: 1.5;">'''
                            section += RevealJSGenerator._format_inline(item.get('segments', []))
                            section += '''</blockquote>'''
                        
                        elif item_type == 'code':
                            # 代码块
                            section += '''
                <pre style="background: #f5f5f5; padding: 10px; border-radius: 5px; margin: 10px 0;"><code>'''
                            section += RevealJSGenerator._escape_html(item.get('text', ''))
                            section += '''</code></pre>'''
                        
                        elif item_type == 'paragraph':
                            # 段落
                            section += '''
                <p style="line-height: 1.5; margin: 10px 0;">'''
                            section += RevealJSGenerator._format_inline(item.get('segments', []))
                            section += '''</p>'''
                        
                        else:  # list
                            # 列表项
                            if not in_list:
                                section += '''
                <ul style="line-height: 1.5; padding-left: 20px; list-style-position: inside;">'''
                                in_list = True
                            
                            section += '''
                    <li>'''
                            section += RevealJSGenerator._format_inline(item.get('segments', []))
                            section += '''</li>'''
                    
                    # 🔧 修复：关闭最后可能打开的 ul
                    if in_list:
                        section += '''
                </ul>'''
                
                section += '''
            </section>'''
            
            sections.append(section)
        
        # 生成主题 CSS
        theme_css = theme.get_reveal_css() if theme else ""
        
        html = f'''<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="utf-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{RevealJSGenerator._escape_html(title)}</title>
    <link rel="stylesheet" href="https://cdn.bootcdn.net/ajax/libs/reveal.js/4.5.0/reset.min.css">
    <link rel="stylesheet" href="https://cdn.bootcdn.net/ajax/libs/reveal.js/4.5.0/reveal.min.css">
    <link rel="stylesheet" href="https://cdn.bootcdn.net/ajax/libs/reveal.js/4.5.0/theme/black.min.css">
    <style>
        {theme_css}
        .reveal ul {{ font-size: 0.8em; }}
        .reveal li {{ margin: 0.5em 0; }}
        .reveal section {{ text-align: left; }}
        .reveal .center {{ text-align: center; }}
    </style>
</head>
<body>
    <div class="reveal">
        <div class="slides">
            {''.join(sections)}
        </div>
    </div>

    <script src="https://cdn.bootcdn.net/ajax/libs/reveal.js/4.5.0/reveal.min.js"></script>
    <script>
        Reveal.initialize({{
            hash: true,
            slideNumber: true,
            transition: 'slide',
            backgroundTransition: 'fade',
            center: false,
            width: 1280,
            height: 720,
            margin: 0.1,
            minScale: 0.2,
            maxScale: 2.0
        }});
    </script>
</body>
</html>'''
        return html


class PowerPointGenerator:
    """PowerPoint 生成器"""
    
    def __init__(self, theme: Theme = None):
        """初始化生成器"""
        if theme is None and get_theme:
            theme = get_theme('default')
        
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(5.625)
        
        # 使用主题颜色
        if theme:
            self.bg_color = theme.get_bg_rgb()
            self.text_color = theme.get_text_rgb()
            self.primary_color = theme.get_primary_rgb()
        else:
            self.bg_color = RGBColor(28, 28, 28)
            self.text_color = RGBColor(255, 255, 255)
            self.primary_color = RGBColor(255, 255, 255)
    
    def add_title_slide(self, title: str, subtitle: str = ""):
        """添加标题幻灯片"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.bg_color
        
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = self.text_color
        
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.alignment = PP_ALIGN.CENTER
            subtitle_para.font.size = Pt(20)
            subtitle_para.font.color.rgb = self.text_color
    
    def add_content_slide(self, title: str, content_items: List[dict]):
        """添加内容幻灯片"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.bg_color
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = self.primary_color
        
        if content_items:
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(4.2))
            text_frame = content_box.text_frame
            text_frame.word_wrap = True
            
            para_index = 0
            for item in content_items:
                item_type = item.get('type', 'list')
                
                if para_index > 0:
                    text_frame.add_paragraph()
                
                p = text_frame.paragraphs[para_index]
                p.font.size = Pt(14)
                p.font.color.rgb = self.text_color
                p.space_before = Pt(4)
                p.line_spacing = 1.5  # 🔧 行间距 1.5
                
                # 根据类型设置样式
                if item_type == 'quote':
                    p.level = 1
                    p.font.italic = True
                elif item_type == 'code':
                    p.font.name = 'Courier New'
                    p.font.size = Pt(12)
                else:
                    p.level = 0
                
                # 添加格式化文本
                segments = item.get('segments', [])
                if item_type == 'code':
                    # 代码块：纯文本
                    p.text = item.get('text', '')
                else:
                    # 其他：格式化文本
                    for seg_idx, seg in enumerate(segments):
                        if seg_idx == 0:
                            p.text = seg.get('text', '')
                            run = p.runs[0]
                        else:
                            run = p.add_run()
                            run.text = seg.get('text', '')
                        
                        run.font.bold = seg.get('bold', False)
                        run.font.italic = seg.get('italic', False)
                        if seg.get('code'):
                            run.font.name = 'Courier New'
                            run.font.size = Pt(12)
                
                para_index += 1
    
    def generate(self, slides_data: List[dict]) -> int:
        """生成 PowerPoint"""
        for slide_data in slides_data:
            if slide_data['type'] == 'title':
                self.add_title_slide(slide_data['title'], slide_data.get('subtitle', ''))
            else:
                self.add_content_slide(slide_data['title'], slide_data['content'])
        
        return len(self.prs.slides)
    
    def save(self, output_path: str):
        """保存 PowerPoint 文件"""
        self.prs.save(output_path)


def convert_markdown_to_presentation(
    markdown_content: Optional[str] = None,
    markdown_file: Optional[str] = None,
    output_path: str = "presentation",
    output_format: str = "html",
    title: Optional[str] = None,
    theme: str = "default"
) -> Union[str, Tuple[str, str]]:
    """
    将 Markdown 转换为演示文稿
    
    Args:
        markdown_content: Markdown 文本内容
        markdown_file: Markdown 文件路径
        output_path: 输出文件路径
        output_format: 输出格式 ("html", "ppt", "both")
        title: 演示文稿标题
        theme: 主题名称
        
    Returns:
        输出文件的绝对路径
    """
    # 获取 Markdown 内容
    if markdown_file:
        with open(markdown_file, 'r', encoding='utf-8') as f:
            content = f.read()
    elif markdown_content:
        content = markdown_content
    else:
        raise ValueError("必须提供 markdown_content 或 markdown_file 参数")
    
    # 如果指定了标题且内容不以一级标题开头，添加标题
    if title and not content.strip().startswith('# '):
        content = f"# {title}\n\n{content}"
    
    # 解析 Markdown
    parser = EnhancedMarkdownParser()
    slides_data = parser.parse(content)
    
    if not slides_data:
        raise ValueError("Markdown 内容为空或格式不正确")
    
    # 智能分页
    original_count = len(slides_data)
    slides_data = SmartPaginator.auto_paginate(slides_data)
    paginated_count = len(slides_data)
    
    if paginated_count > original_count:
        print(f"初始解析: 共 {original_count} 张幻灯片")
        print(f"智能分页后: 共 {paginated_count} 张幻灯片")
    
    # 确定标题
    if not title and slides_data and slides_data[0]['type'] == 'title':
        title = slides_data[0]['title']
    elif not title:
        title = "演示文稿"
    
    # 获取主题
    theme_obj = None
    if get_theme:
        theme_obj = get_theme(theme)
        print(f"使用主题: {theme}")
    
    # 处理输出路径
    output_path = Path(output_path)
    
    # 生成输出
    results = []
    
    if output_format in ("html", "both"):
        # 生成 HTML
        html_generator = RevealJSGenerator()
        html_content = html_generator.generate(slides_data, title, theme_obj)
        
        # 确定 HTML 输出路径
        if output_format == "both":
            html_path = output_path.with_suffix('.html')
        elif output_path.suffix == '.html':
            html_path = output_path
        else:
            html_path = output_path.with_suffix('.html')
        
        html_path = html_path.resolve()
        html_path.write_text(html_content, encoding='utf-8')
        
        print(f"\nHTML 演示文稿创建成功！")
        print(f"文件保存位置: {html_path}")
        print(f"共 {len(slides_data)} 张幻灯片")
        print(f"使用方法: 在浏览器中打开，使用方向键或空格键翻页")
        
        results.append(str(html_path))
    
    if output_format in ("ppt", "both"):
        # 生成 PowerPoint
        ppt_generator = PowerPointGenerator(theme_obj)
        slide_count = ppt_generator.generate(slides_data)
        
        # 确定 PPT 输出路径
        if output_format == "both":
            ppt_path = output_path.with_suffix('.pptx')
        elif output_path.suffix == '.pptx':
            ppt_path = output_path
        else:
            ppt_path = output_path.with_suffix('.pptx')
        
        ppt_path = ppt_path.resolve()
        ppt_generator.save(str(ppt_path))
        
        print(f"\nPowerPoint 演示文稿创建成功！")
        print(f"文件保存位置: {ppt_path}")
        print(f"共 {slide_count} 张幻灯片")
        
        results.append(str(ppt_path))
    
    # 返回结果
    if output_format == "both":
        return tuple(results)
    else:
        return results[0]


def main():
    """命令行入口"""
    parser = argparse.ArgumentParser(
        description="将 Markdown 文档转换为演示文稿（HTML 或 PowerPoint）"
    )
    parser.add_argument(
        '-i', '--input',
        required=True,
        help="输入 Markdown 文件路径"
    )
    parser.add_argument(
        '-o', '--output',
        default="presentation",
        help="输出文件路径（默认: presentation）"
    )
    parser.add_argument(
        '-f', '--format',
        choices=['html', 'ppt', 'both'],
        default='html',
        help="输出格式: html（reveal.js）, ppt（PowerPoint）, both（两者都生成）"
    )
    parser.add_argument(
        '-t', '--title',
        help="演示文稿标题（可选）"
    )
    parser.add_argument(
        '--theme',
        choices=['philosophy', 'business', 'minimal', 'traditional', 'dark', 'default'],
        default='default',
        help="演示文稿主题"
    )
    parser.add_argument(
        '--list-themes',
        action='store_true',
        help="列出所有可用主题"
    )
    
    args = parser.parse_args()
    
    # 列出主题
    if args.list_themes:
        if list_themes:
            list_themes()
        else:
            print("主题模块未加载")
        sys.exit(0)
    
    try:
        result = convert_markdown_to_presentation(
            markdown_file=args.input,
            output_path=args.output,
            output_format=args.format,
            title=args.title,
            theme=args.theme
        )
        sys.exit(0)
    except Exception as e:
        print(f"错误: {e}", file=sys.stderr)
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()

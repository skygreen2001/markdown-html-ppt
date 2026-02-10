#!/usr/bin/env python3
"""
Markdown Html PowerPoint Converter
将 Markdown 文档转换为 PowerPoint 演示文稿
"""

import re
import argparse
import sys
from pathlib import Path
from typing import List, Tuple, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


class MarkdownToPPTConverter:
    """Markdown 到 PowerPoint 转换器"""
    
    def __init__(self, bg_color=(28, 28, 28), text_color=(255, 255, 255)):
        """
        初始化转换器
        
        Args:
            bg_color: 背景颜色 RGB 元组
            text_color: 文字颜色 RGB 元组
        """
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(5.625)
        self.bg_color = RGBColor(*bg_color)
        self.text_color = RGBColor(*text_color)
        
    def parse_markdown(self, content: str) -> List[dict]:
        """
        解析 Markdown 内容为结构化数据
        
        Args:
            content: Markdown 文本内容
            
        Returns:
            幻灯片数据列表
        """
        slides = []
        current_slide = None
        lines = content.split('\n')
        
        for line in lines:
            line = line.rstrip()
            
            # 一级标题 - 标题幻灯片
            if line.startswith('# '):
                if current_slide:
                    slides.append(current_slide)
                current_slide = {
                    'type': 'title',
                    'title': line[2:].strip(),
                    'subtitle': '',
                    'content': []
                }
            
            # 二级/三级标题 - 内容幻灯片
            elif line.startswith('## ') or line.startswith('### '):
                if current_slide:
                    slides.append(current_slide)
                level = 2 if line.startswith('## ') else 3
                title = line[level+1:].strip()
                current_slide = {
                    'type': 'content',
                    'title': title,
                    'content': []
                }
            
            # 列表项
            elif line.strip().startswith(('- ', '* ', '+ ')) or re.match(r'^\s*\d+\.\s', line):
                if current_slide:
                    # 移除列表标记
                    content_text = re.sub(r'^[\s\-\*\+\d\.]+\s*', '', line).strip()
                    if content_text:
                        current_slide['content'].append(content_text)
            
            # 普通段落
            elif line.strip() and current_slide:
                # 如果是标题幻灯片且还没有副标题，作为副标题
                if current_slide['type'] == 'title' and not current_slide['subtitle']:
                    current_slide['subtitle'] = line.strip()
                # 否则作为内容
                elif line.strip():
                    current_slide['content'].append(line.strip())
        
        # 添加最后一张幻灯片
        if current_slide:
            slides.append(current_slide)
        
        return slides
    
    def add_title_slide(self, title: str, subtitle: str = ""):
        """添加标题幻灯片"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 设置背景
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.bg_color
        
        # 添加标题
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = self.text_color
        
        # 添加副标题
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.alignment = PP_ALIGN.CENTER
            subtitle_para.font.size = Pt(20)
            subtitle_para.font.color.rgb = self.text_color
    
    def add_content_slide(self, title: str, content_list: List[str]):
        """添加内容幻灯片"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 设置背景
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.bg_color
        
        # 添加标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = self.text_color
        
        # 添加内容
        if content_list:
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(4.2))
            text_frame = content_box.text_frame
            text_frame.word_wrap = True
            
            for i, item in enumerate(content_list):
                if i > 0:
                    text_frame.add_paragraph()
                p = text_frame.paragraphs[i]
                p.text = item
                p.font.size = Pt(14)
                p.font.color.rgb = self.text_color
                p.space_before = Pt(4)
                p.level = 0
    
    def convert(self, markdown_content: str) -> int:
        """
        执行转换
        
        Args:
            markdown_content: Markdown 文本内容
            
        Returns:
            生成的幻灯片数量
        """
        slides_data = self.parse_markdown(markdown_content)
        
        for slide_data in slides_data:
            if slide_data['type'] == 'title':
                self.add_title_slide(slide_data['title'], slide_data.get('subtitle', ''))
            else:
                self.add_content_slide(slide_data['title'], slide_data['content'])
        
        return len(self.prs.slides)
    
    def save(self, output_path: str):
        """
        保存 PowerPoint 文件
        
        Args:
            output_path: 输出文件路径
        """
        self.prs.save(output_path)


def convert_markdown_to_ppt(
    markdown_content: Optional[str] = None,
    markdown_file: Optional[str] = None,
    output_path: str = "presentation.pptx",
    title: Optional[str] = None
) -> str:
    """
    将 Markdown 转换为 PowerPoint
    
    Args:
        markdown_content: Markdown 文本内容（与 markdown_file 二选一）
        markdown_file: Markdown 文件路径（与 markdown_content 二选一）
        output_path: 输出文件路径
        title: 演示文稿标题（可选，会添加到第一张幻灯片）
        
    Returns:
        输出文件的绝对路径
    """
    # 获取 Markdown 内容
    if markdown_file:
        with open(markdown_file, 'r', encoding='utf-8') as f:
            content = f.read()
    elif markdown_content:
        content = markdown_content
    else:
        raise ValueError("必须提供 markdown_content 或 markdown_file 参数")
    
    # 如果指定了标题且内容不以一级标题开头，添加标题
    if title and not content.strip().startswith('# '):
        content = f"# {title}\n\n{content}"
    
    # 创建转换器并执行转换
    converter = MarkdownToPPTConverter()
    slide_count = converter.convert(content)
    
    # 保存文件
    output_path = str(Path(output_path).resolve())
    converter.save(output_path)
    
    print(f"PowerPoint 演示文稿创建成功！")
    print(f"文件保存位置: {output_path}")
    print(f"共 {slide_count} 张幻灯片")
    
    return output_path


def main():
    """命令行入口"""
    parser = argparse.ArgumentParser(
        description="将 Markdown 文档转换为 PowerPoint 演示文稿"
    )
    parser.add_argument(
        '-i', '--input',
        required=True,
        help="输入 Markdown 文件路径"
    )
    parser.add_argument(
        '-o', '--output',
        default="presentation.pptx",
        help="输出 PowerPoint 文件路径（默认: presentation.pptx）"
    )
    parser.add_argument(
        '-t', '--title',
        help="演示文稿标题（可选）"
    )
    
    args = parser.parse_args()
    
    try:
        output_path = convert_markdown_to_ppt(
            markdown_file=args.input,
            output_path=args.output,
            title=args.title
        )
        sys.exit(0)
    except Exception as e:
        print(f"错误: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
